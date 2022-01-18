# 2개 넣기

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

#insert 슬라이드
blank_slide = prs.slide_layouts[6]
slide_1 = prs.slides.add_slide(blank_slide)

#기본 위치 및 크기
left = Cm(0.85)
top = Cm(1.67)     # left , top은 상대 위치
width = Cm(32.14)     # 너비, 높이는 텍스트 상자 크기
height = Cm(6.63)

# 지정된 위치에 텍스트 상자 추가
textbox = slide_1.shapes.add_textbox (left, top, width, height)
tf = textbox.text_frame

para = tf.add_paragraph()     # 단락
para.text = 'My word'   # 단락에 텍스트
para.alignment = PP_ALIGN.CENTER
para.line_spacing = 1.5 # 1.5배 줄 간격

# ## 글꼴
font = para.font
font.name = '함초름돋음'     #글꼴 유형 font.bold
font.bold = True     #bold font.size
font.size = Pt(50)     #크기

# 새 텍스트 상자
new_para = tf.add_paragraph()  # 添加段落
new_para.text = '\t\t\t元 · 管道昇 '  # 段落文字
new_para.alignment = PP_ALIGN.CENTER    # 居中
new_para.line_spacing = 1.5

# ## 글꼴
font = new_para.font
font.name = '함초름돋음'     #글꼴 유형 font.size
font.size = Pt(50)     #size font.underline
font.underline = True     #underline

prs.save('test_slide.pptx')