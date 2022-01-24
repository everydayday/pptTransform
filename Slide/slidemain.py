# from get_lyrics import get_lyrics
from math import ceil  # slide 갯수구하기 by 올림

import win32com.client
from pptx import Presentation  # 라이브러리
from pptx.util import Pt, Cm


# from pptx_tools import utils


def get_slide_num(lyric_str):
    slide_num = ceil(len(lyric_str) / 2)
    print("here is get_slide_num def")
    print(lyric_str.count('\n'))
    return slide_num


# slide_str 만들기 / index : 슬라이드마다 들어가는 가사 두줄
def get_slide_str(lyric_str):
    my_list = lyric_str.split('\n')
    slide_list = []
    idx_count = 0

    for idx, value in enumerate(my_list):
        if idx % 2 == 0 and value == '':
            my_list[idx] = '\n'
            my_list.insert(idx + 1, '\n')

    print("here is my_list", my_list)
    for idx, value in enumerate(my_list):

        if idx % 2 == 0:  # 괜히 try문 안 쓰는게 나을 듯.
            if value == '':
                slide_list.append('\n')
                slide_list.append('\n')
                idx_count += 1
            else:
                slide_list.append(value)
        else:
            if value == '\n':
                slide_list[idx_count] = slide_list[idx_count] + value
                idx_count += 1
            else:
                slide_list[idx_count] = slide_list[idx_count] + "\n" + value
                idx_count += 1

    print("here is slide list \n", slide_list)
    return slide_list
    slide_str = []
    slide_num = get_slide_num(lyric_str)
    '''
    for idx in range(lyric_str):
        slide_str[idx // 2] = slide_str[idx // 2] + value
    '''
    for i in range(slide_num):
        try:
            slide_str.append(lyric_str[i * 2] + lyric_str[i * 2 + 1].replace("\n", ""))
        except IndexError:  # 가사줄의 갯수가 홀수여서 마지막에 str[i*2]까지만 해당할 시
            slide_str.append(lyric_str[i * 2])

    return slide_str


# slide 생성
# 가사만 주면 ok
def make_slide(lyric_str, font_size=30, font_style='함초름돋음'):  # slide_str 아무 값도 없을 시, 오류 메시지 출력
    # 파워포인트 객체 선언
    prs = Presentation()
    # slide 마다 가사 리스트
    slide_str = get_slide_str(lyric_str)
    # slide 갯수 구하기 (올림)
    global slide_num
    slide_num = len(slide_str)
    # slide_num = ceil(len(lyric_str) / 2)
    # slide 만들기
    for i in range(0, slide_num):
        blank_slide_layout = prs.slide_layouts[6]  # 슬라이드 종류 선택
        slide = prs.slides.add_slide(blank_slide_layout)  # 슬라이드 추가

        '''left = Inches(1)

        top = Inches(1)

        width = Inches(5)

        height = Inches(0.5)'''
        # 위치, 가로/세로 길이
        left = Cm(1.91)

        top = Cm(7.46)
        width = Cm(4.08)
        height = Cm(21.59)

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.text = slide_str[i]
        tf.paragraphs[0].font.name = font_style
        tf.paragraphs[0].font.size = Pt(font_size)

        # title = prs.slides[i].shapes.title
        # content = slide.placeholders[0]
        # content.text = slide_str[i]
        # p = title.text_frame.paragraphs[0]
        # run = p.add_run()

        # title.text_frame.paragraphs[0].font.name = font_style
        print(tf.text)
        print("----------------------------")
        # title.text_frame.paragraphs[1].font.name = font_style #paragraph : 단락/행/엔터기준 구분
        # title.text_frame.paragraphs[0].font.size = Pt(font_size)
        # title.text_frame.paragraphs[1].font.size = Pt(font_size)

    print("before prs.save")
    prs.save('add all slides1.pptx')


# image 저장
def save_pptx_as_png():
    '''if os.path.isdir(png_foldername) and not overwrite_folder:
        print(f"Folder {png_foldername} already exists. "
              f"Set overwrite_folder=True, if you want to overwrite folder content.")
        return

    powerpoint = CreateObject("Powerpoint.Application")
    pp_constants = Constants(powerpoint)

    pres = powerpoint.Presentations.Open(pptx_filename)
    pres.SaveAs(png_foldername, pp_constants.ppSaveAsPNG)
    pres.close()
    if powerpoint.Presentations.Count == 0:  # only close, when no other Presentations are open!
        powerpoint.quit()'''
    Application = win32com.client.Dispatch("PowerPoint.Application")
    Presentation = Application.Presentations.Open(r"C:\Users\김대희\PycharmProjects\pythonProject2\UI\add all slides1.pptx")
    i = 0

    print("before global slide_num")

    print("number of slide_num",slide_num)
    '''
    for i in range(slide_num):
        Presentation.Slides[i].Export(r"C:Users\김대희\PycharmProjects\pythonProject2\slideImagefolder\pptimage{}.jpg".format(i), "JPG")
    '''

    #Presentation.Slides[0].Export(r"C:\Users\김대희\PycharmProjects\pythonProject2\slideImagefolder\pptimage{}.jpg".format(0), "JPG")

    Application.Quit()
    Presentation = None
    Application = None


def slide_main(lyrics):
    # 폰트 크기, 스타일 입력
    while True:
        try:  # 잘못입력했을 시
            size = 30
            style = "함초름돋음"

            if size == '':  # 디폴트 값
                font_size = 30
            else:  # 주어진 값 있을 시
                font_size = int(size)
            if style == '':  # 디폴트 값
                font_style = '함초름돋음'
            else:  ## 이상한 폰트 입력해도 입력이 된다는 문제점 ##
                font_style = style
            break  # 오류 없으면 while문 넘어가기
        except ValueError:  # 숫자가 아닌 값 입력 시
            print("잘못된 값을 입력했습니다.")
            print("처음부터 다시 입력하세요.")

    lyric_str = lyrics
    make_slide(lyric_str, font_size, font_style)
    save_pptx_as_png()
    # png_folder = 'C:/Users/김대희/PycharmProjects/pythonProject2/Slide'
    # pptx_file = 'C:/Users/김대희/PycharmProjects/pythonProject2/Slide/add all slides1.pptx'
    # utils.save_pptx_as_png(png_folder, pptx_file)


if __name__ == '__main__':
    slide_main()
