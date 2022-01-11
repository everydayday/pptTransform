from pptx import Presentation  # 라이브러리
from pptx.util import Pt
from get_lyrics import get_lyrics
from math import ceil # slide 갯수구하기 by 올림

# 전체 가사 입력
str = get_lyrics()
slide_num = ceil(len(str)/2)   # slide 갯수 (올림)


slide_str = [] # index마다 입력되는 2줄의 가사 존재

# slide_str 만들기
for i in range(slide_num):
    try:
        slide_str.append(str[i*2] + '\n' + str[i*2+1])
    except IndexError: # 가사줄의 갯수가 홀수여서 str[i*2]까지만 해당할 시
        slide_str.append(str[i*2] + '\n')



# 파워포인트 객체 선언
prs = Presentation()

# 슬라이드 생성

for i in range(0, slide_num):
    title_slide_layout = prs.slide_layouts[0]  # 슬라이드 종류 선택
    slide = prs.slides.add_slide(title_slide_layout)  # 슬라이드 추가
    content = slide.placeholders[0]
    content.text = slide_str[i]
    title = prs.slides[i].shapes.title
    title.text_frame.paragraphs[0].font.name = '함초름돋음'
    title.text_frame.paragraphs[1].font.name = '함초름돋음'
    title.text_frame.paragraphs[0].font.size = Pt(30)
    title.text_frame.paragraphs[1].font.size = Pt(30)



prs.save('add all slides1.pptx')