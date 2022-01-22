from pptx import Presentation

from pptx.util import Inches, Pt

prs = Presentation()

blank_slide_layout = prs.slide_layouts[6]

slide = prs.slides.add_slide(blank_slide_layout)

# 위치 좌표, 가로/세로 길이

left = Inches(1)

top = Inches(1)

width = Inches(5)

height = Inches(0.5)

tb = slide.shapes.add_textbox(left, top, width, height)

tf = tb.text_frame

tf.text = 'Text Box'

p = tf.add_paragraph()

p.text = "This is a second paragraph that's bold"

p.font.bold = True

p = tf.add_paragraph()

p.text = "This is a third paragraph that's big"

p.font.size = Pt(40)

prs.save('demo.pptx')